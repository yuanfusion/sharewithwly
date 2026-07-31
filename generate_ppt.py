# -*- coding: utf-8 -*-
"""物理AI主题PPT生成脚本 —— 20页，清新浅色风格"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------------- 配色（清新浅色系） ----------------
INK      = RGBColor(0x0F, 0x2B, 0x2A)   # 深墨绿黑（标题）
BODY     = RGBColor(0x3E, 0x50, 0x4F)   # 正文灰绿
MUTED    = RGBColor(0x8A, 0x9B, 0x99)   # 弱化灰
TEAL     = RGBColor(0x0D, 0x94, 0x88)   # 主色 青碧
TEAL_DK  = RGBColor(0x0B, 0x6E, 0x64)
MINT     = RGBColor(0xCC, 0xFB, 0xF1)   # 浅薄荷
MINT_BG  = RGBColor(0xF0, 0xFD, 0xFA)   # 卡片底色 薄荷白
SKY      = RGBColor(0x02, 0x8A, 0xC8)   # 晴空蓝
SKY_BG   = RGBColor(0xF0, 0xF9, 0xFF)
AMBER    = RGBColor(0xD9, 0x7D, 0x0E)   # 暖橙（点缀）
AMBER_BG = RGBColor(0xFF, 0xFB, 0xEB)
ROSE     = RGBColor(0xE0, 0x4E, 0x6A)
ROSE_BG  = RGBColor(0xFF, 0xF5, 0xF6)
VIOLET   = RGBColor(0x6D, 0x5F, 0xC4)
VIOLET_BG= RGBColor(0xF5, 0xF4, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xFC, 0xFE, 0xFD)   # 页面底色：极浅暖白
LINE     = RGBColor(0xDD, 0xEA, 0xE7)   # 浅描边
SOFT     = RGBColor(0xF4, 0xF8, 0xF7)   # 浅灰绿卡片

FONT = "微软雅黑"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ---------------- 基础工具 ----------------
def set_run_font(run, size, color, bold=False, name=FONT, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        for e in rPr.findall(qn(tag)):
            rPr.remove(e)
        e = rPr.makeelement(qn(tag), {"typeface": name})
        rPr.append(e)

def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, space_after=0, wrap=True):
    """runs: list[ list[(text,size,color,bold)] ] 每个子列表为一段"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for (t, s, c, b) in para:
            r = p.add_run()
            r.text = t
            set_run_font(r, s, c, b)
    return tb

def add_shape(slide, shape, x, y, w, h, fill=None, line=None, line_w=0.75,
              shadow=False, adj=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if adj is not None:
        try:
            sp.adjustments[0] = adj
        except Exception:
            pass
    return sp

def card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=0.08, line_w=1.0):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                     fill=fill, line=line, line_w=line_w, adj=radius)

def hline(slide, x, y, w, color=LINE, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H, fill=BG)
    return s

def footer(slide, page, total=20):
    add_text(slide, MARGIN, Inches(7.08), Inches(4), Inches(0.3),
             [[("物理 AI · Physical AI", 9, MUTED, False)]])
    add_text(slide, EMU_W - MARGIN - Inches(1.2), Inches(7.08), Inches(1.2), Inches(0.3),
             [[(f"{page:02d} / {total}", 9, MUTED, False)]], align=PP_ALIGN.RIGHT)

def header(slide, kicker, title, page):
    # 顶部小节标 + 标题 + 装饰
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(0.52), Inches(0.5), Inches(0.14),
              fill=TEAL, adj=0.5)
    add_text(slide, MARGIN + Inches(0.62), Inches(0.4), Inches(9), Inches(0.35),
             [[(kicker, 12, TEAL, True)]])
    add_text(slide, MARGIN, Inches(0.74), Inches(11.5), Inches(0.6),
             [[(title, 26, INK, True)]])
    footer(slide, page)

def deco_circles(slide, right=True):
    """清新背景装饰：浅色圆点/圆环"""
    if right:
        add_shape(slide, MSO_SHAPE.OVAL, Inches(11.6), Inches(0.35), Inches(0.28), Inches(0.28), fill=MINT)
        add_shape(slide, MSO_SHAPE.OVAL, Inches(12.05), Inches(0.55), Inches(0.16), Inches(0.16), fill=SKY_BG, line=SKY, line_w=1.0)
        add_shape(slide, MSO_SHAPE.OVAL, Inches(12.35), Inches(0.3), Inches(0.1), Inches(0.1), fill=TEAL)

def chip(slide, x, y, w, text, fg=TEAL_DK, bg=MINT):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34), fill=bg, adj=0.5)
    add_text(slide, x, y + Inches(0.045), w, Inches(0.26),
             [[(text, 11, fg, True)]], align=PP_ALIGN.CENTER)

def section_slide(no, zh, en, points, page):
    s = new_slide()
    # 左侧大色块装饰（浅色）
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-1.2), Inches(-1.0), Inches(6.2), Inches(9.5),
              fill=MINT_BG, adj=0.12)
    add_shape(s, MSO_SHAPE.OVAL, Inches(3.4), Inches(4.9), Inches(1.6), Inches(1.6), fill=MINT)
    add_shape(s, MSO_SHAPE.OVAL, Inches(3.75), Inches(5.25), Inches(0.9), Inches(0.9), fill=WHITE)
    add_text(s, Inches(0.9), Inches(2.0), Inches(4), Inches(1.6),
             [[(no, 100, TEAL, True)]])
    add_text(s, Inches(0.95), Inches(3.78), Inches(5.1), Inches(0.4),
             [[(en, 11.5, MUTED, True)]], wrap=False)
    # 右侧标题
    add_text(s, Inches(5.6), Inches(2.35), Inches(7.0), Inches(1.0),
             [[(zh, 40, INK, True)]])
    hline(s, Inches(5.62), Inches(3.3), Inches(6.9), color=LINE, weight=1.2)
    yy = 3.62
    for p in points:
        add_shape(s, MSO_SHAPE.OVAL, Inches(5.66), Inches(yy + 0.09), Inches(0.1), Inches(0.1), fill=TEAL)
        add_text(s, Inches(5.92), Inches(yy), Inches(6.6), Inches(0.4),
                 [[(p, 15, BODY, False)]])
        yy += 0.52
    footer(s, page)
    return s

# =====================================================================
# P1 封面
# =====================================================================
s = new_slide()
# 背景装饰：浅色圆弧与圆点
add_shape(s, MSO_SHAPE.OVAL, Inches(9.0), Inches(-2.6), Inches(7.2), Inches(7.2), fill=MINT_BG)
add_shape(s, MSO_SHAPE.OVAL, Inches(10.1), Inches(-1.6), Inches(5.0), Inches(5.0), fill=MINT)
add_shape(s, MSO_SHAPE.OVAL, Inches(11.0), Inches(4.9), Inches(3.4), Inches(3.4), fill=SKY_BG)
add_shape(s, MSO_SHAPE.OVAL, Inches(0.5), Inches(6.3), Inches(0.18), Inches(0.18), fill=TEAL)
add_shape(s, MSO_SHAPE.OVAL, Inches(0.95), Inches(6.3), Inches(0.18), Inches(0.18), fill=SKY)
add_shape(s, MSO_SHAPE.OVAL, Inches(1.4), Inches(6.3), Inches(0.18), Inches(0.18), fill=AMBER)

chip(s, Inches(0.95), Inches(1.5), Inches(2.6), "行业科普 · 趋势解读", fg=TEAL_DK, bg=MINT)
add_text(s, Inches(0.92), Inches(2.15), Inches(11.4), Inches(2.2),
         [[("物理 ", 66, INK, True), ("AI", 66, TEAL, True)],
          [("Physical AI：让智能走出屏幕，走进真实世界", 24, TEAL_DK, True)]],
         line_spacing=1.25)
add_text(s, Inches(0.95), Inches(4.55), Inches(9.6), Inches(0.9),
         [[("从会“说”的 AI，到会“做”的 AI —— 机器人与自动驾驶背后的下一代人工智能", 15, BODY, False)]])
hline(s, Inches(0.97), Inches(5.35), Inches(4.2), color=LINE, weight=1.5)
add_text(s, Inches(0.95), Inches(5.6), Inches(9), Inches(0.4),
         [[("内部技术分享  |  2026 年 7 月", 13, MUTED, False)]])

# 右侧“感知-理解-行动”小图示
cx = Inches(10.55)
for i, (t, c, bgc) in enumerate([("感知", TEAL, WHITE), ("理解", SKY, WHITE), ("行动", AMBER, WHITE)]):
    yy = Inches(2.6 + i * 0.78)
    add_shape(s, MSO_SHAPE.OVAL, cx, yy, Inches(0.62), Inches(0.62), fill=bgc, line=c, line_w=1.5)
    add_text(s, cx, yy + Inches(0.16), Inches(0.62), Inches(0.3),
             [[(t, 12, c, True)]], align=PP_ALIGN.CENTER)
    if i < 2:
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx + Inches(0.31), yy + Inches(0.62),
                                    cx + Inches(0.31), yy + Inches(0.78))
        ln.line.color.rgb = MUTED; ln.line.width = Pt(1.2); ln.shadow.inherit = False
add_text(s, Inches(11.4), Inches(2.6), Inches(1.5), Inches(2.4),
         [[("与真实世界", 13, BODY, True)], [("持续闭环交互", 13, BODY, True)]], line_spacing=1.4)

# =====================================================================
# P2 目录
# =====================================================================
s = new_slide()
header(s, "CONTENTS", "目录", 2)
deco_circles(s)
items = [
    ("01", "什么是物理 AI", "定义、演进与核心能力", TEAL, MINT_BG),
    ("02", "物理 AI 在做什么", "三大核心任务与技术底座", SKY, SKY_BG),
    ("03", "与 Agent、机器人的对比", "特点与优势在哪里", VIOLET, VIOLET_BG),
    ("04", "产业格局与优秀案例", "领先公司 + 落地故事", AMBER, AMBER_BG),
    ("05", "未来趋势与挑战", "五大趋势、四大挑战", ROSE, ROSE_BG),
    ("06", "普通人如何入门", "学习路线与资源清单", TEAL_DK, MINT_BG),
]
for i, (no, t, d, c, bgc) in enumerate(items):
    col, row = i % 2, i // 2
    x = MARGIN + col * Inches(6.15)
    y = Inches(1.75) + row * Inches(1.72)
    card(s, x, y, Inches(5.85), Inches(1.5), fill=bgc, line=None)
    add_text(s, x + Inches(0.32), y + Inches(0.3), Inches(1.1), Inches(0.9),
             [[(no, 30, c, True)]])
    add_text(s, x + Inches(1.45), y + Inches(0.32), Inches(4.2), Inches(0.5),
             [[(t, 18, INK, True)]])
    add_text(s, x + Inches(1.45), y + Inches(0.86), Inches(4.2), Inches(0.4),
             [[(d, 12, BODY, False)]])

# =====================================================================
# P3 章节页 01
# =====================================================================
section_slide("01", "什么是物理 AI", "WHAT IS PHYSICAL AI",
              ["AI 的四次范式演进：从感知 AI 到物理 AI",
               "物理 AI 的定义与三大核心特征",
               "为什么说物理 AI 的“ChatGPT 时刻”已经到来"], 3)

# =====================================================================
# P4 AI 的四次范式演进（时间线）
# =====================================================================
s = new_slide()
header(s, "AI EVOLUTION", "AI 的四次范式演进", 4)
deco_circles(s)
add_text(s, MARGIN, Inches(1.42), Inches(12), Inches(0.4),
         [[("英伟达 CEO 黄仁勋将 AI 的发展划分为四个阶段，物理 AI 是最新的一波浪潮", 13, BODY, False)]])
stages = [
    ("感知 AI", "Perception AI", "约 2012 年起", "看得懂、听得清", "图像识别、语音转写、人脸识别", MUTED, SOFT, False),
    ("生成式 AI", "Generative AI", "2022 年起", "会写作、会画画", "ChatGPT、文生图、代码生成", SKY, SKY_BG, False),
    ("代理 AI", "Agentic AI", "2024 年起", "会推理、会调用工具", "AI Agent 自动完成数字任务", VIOLET, VIOLET_BG, False),
    ("物理 AI", "Physical AI", "2025 年起", "理解物理规律、能动手", "机器人、自动驾驶走进现实", TEAL, MINT_BG, True),
]
line_y = Inches(3.05)
hline(s, MARGIN + Inches(0.2), line_y, Inches(11.7), color=LINE, weight=2.0)
for i, (t, en, yr, cap, eg, c, bgc, hot) in enumerate(stages):
    x = MARGIN + i * Inches(3.08)
    add_shape(s, MSO_SHAPE.OVAL, x + Inches(1.25), line_y - Inches(0.09), Inches(0.18), Inches(0.18),
              fill=c if not hot else TEAL)
    cw, chh = Inches(2.78), Inches(2.5)
    cy = Inches(3.45)
    card(s, x, cy, cw, chh, fill=bgc, line=(TEAL if hot else None), line_w=1.5)
    if hot:
        chip(s, x + Inches(1.52), cy - Inches(0.17), Inches(1.15), "当下浪潮", fg=WHITE, bg=TEAL)
    add_text(s, x + Inches(0.25), cy + Inches(0.22), cw - Inches(0.5), Inches(0.45),
             [[(t, 19, INK, True)]])
    add_text(s, x + Inches(0.25), cy + Inches(0.66), cw - Inches(0.5), Inches(0.3),
             [[(en + " · " + yr, 10.5, c if not hot else TEAL_DK, True)]])
    add_text(s, x + Inches(0.25), cy + Inches(1.05), cw - Inches(0.5), Inches(0.4),
             [[(cap, 14, BODY, True)]])
    add_text(s, x + Inches(0.25), cy + Inches(1.55), cw - Inches(0.5), Inches(0.8),
             [[(eg, 11.5, MUTED, False)]], line_spacing=1.3)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.3), Inches(12.1), Inches(0.62),
          fill=WHITE, line=LINE, adj=0.5)
add_text(s, MARGIN, Inches(6.44), Inches(12.1), Inches(0.4),
         [[("一句话理解：前几代 AI 生活在“屏幕里”，物理 AI 则让智能第一次真正走进“物理世界”。", 14, TEAL_DK, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P5 物理 AI 的定义
# =====================================================================
s = new_slide()
header(s, "DEFINITION", "什么是物理 AI？", 5)
deco_circles(s)
card(s, MARGIN, Inches(1.5), Inches(12.1), Inches(1.35), fill=MINT_BG, line=None, radius=0.12)
add_text(s, MARGIN + Inches(0.4), Inches(1.72), Inches(11.3), Inches(1.0),
         [[("物理 AI（Physical AI）是能够理解现实世界、并与之交互的 AI 模型 ——", 17, INK, True)],
          [("它让机器人、自动驾驶汽车等自主机器，在真实物理世界中完成 感知 → 理解 → 行动 的闭环。", 15, TEAL_DK, True)]],
         line_spacing=1.35)
feats = [
    ("懂物理规律", "理解重力、摩擦、碰撞、材质特性，知道“松手物体会下落”“水杯一推会洒”", TEAL, MINT_BG),
    ("多模态感知", "通过摄像头、激光雷达、麦克风、触觉传感器，像人一样“眼观六路、耳听八方”", SKY, SKY_BG),
    ("能动手执行", "不只输出文字，而是输出动作：抓取、行走、驾驶，直接改变物理世界", AMBER, AMBER_BG),
    ("从交互中学习", "在仿真与真实环境中反复试错，越用越聪明，能力可以持续进化", VIOLET, VIOLET_BG),
]
for i, (t, d, c, bgc) in enumerate(feats):
    x = MARGIN + (i % 2) * Inches(6.15)
    y = Inches(3.15) + (i // 2) * Inches(1.62)
    card(s, x, y, Inches(5.85), Inches(1.42), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.12), Inches(1.42), fill=c, adj=0.5)
    add_text(s, x + Inches(0.35), y + Inches(0.18), Inches(5.3), Inches(0.4),
             [[(f"特征 {i+1} · {t}", 15, INK, True)]])
    add_text(s, x + Inches(0.35), y + Inches(0.62), Inches(5.3), Inches(0.7),
             [[(d, 12, BODY, False)]], line_spacing=1.25)
add_text(s, MARGIN, Inches(6.55), Inches(12.1), Inches(0.4),
         [[("2026 年 CES 上，黄仁勋宣布：物理 AI 的“ChatGPT 时刻”已经到来。", 13, AMBER, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P6 物理 AI 在做什么（三大核心任务）
# =====================================================================
s = new_slide()
header(s, "WHAT IT DOES", "物理 AI 主要在做什么？", 6)
deco_circles(s)
add_text(s, MARGIN, Inches(1.42), Inches(12), Inches(0.4),
         [[("核心思路：先在“虚拟世界”里低成本学会本领，再把本领迁移到真实机器上", 13, BODY, False)]])
tasks = [
    ("① 建世界", "构建高保真仿真环境与世界模型", ["用数字孪生复刻工厂、道路、家庭",
     "世界模型按物理规律“想象”未来画面", "代表：NVIDIA Omniverse / Cosmos"]),
    ("② 教本领", "在仿真中生成数据、训练模型", ["合成数据解决真实数据稀缺昂贵问题",
     "强化学习 + 模仿学习，一夜可练百万次", "代表：Isaac Lab / GR00T 模型"]),
    ("③ 上真机", "把模型部署到机器人与车辆", ["Sim-to-Real 迁移到真实硬件",
     "端侧算力实时推理，闭环控制", "代表：Jetson Thor / DRIVE 平台"]),
]
for i, (t, st, pts) in enumerate(tasks):
    x = MARGIN + i * Inches(4.13)
    card(s, x, Inches(1.95), Inches(3.85), Inches(3.6), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.95), Inches(3.85), Inches(0.1), fill=TEAL, adj=0.5)
    add_text(s, x + Inches(0.3), Inches(2.25), Inches(3.2), Inches(0.5), [[(t, 20, INK, True)]])
    add_text(s, x + Inches(0.3), Inches(2.78), Inches(3.3), Inches(0.6),
             [[(st, 12.5, TEAL_DK, True)]], line_spacing=1.2)
    yy = 3.5
    for p in pts:
        add_shape(s, MSO_SHAPE.OVAL, x + Inches(0.32), Inches(yy + 0.08), Inches(0.09), Inches(0.09), fill=TEAL)
        add_text(s, x + Inches(0.52), Inches(yy), Inches(3.15), Inches(0.75),
                 [[(p, 11.5, BODY, False)]], line_spacing=1.15)
        yy += 0.68
    if i < 2:
        add_shape(s, MSO_SHAPE.CHEVRON, x + Inches(3.87), Inches(3.5), Inches(0.3), Inches(0.4), fill=MINT)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(5.85), Inches(12.1), Inches(0.95),
          fill=SKY_BG, line=None, adj=0.18)
add_text(s, MARGIN + Inches(0.4), Inches(6.05), Inches(11.5), Inches(0.6),
         [[("为什么重要？  真实世界试错又慢又贵又危险 —— 让机器在仿真里“摔一百万次”，", 13.5, INK, True),
           ("才能把成本、周期与安全风险降到可落地的水平。", 13.5, SKY, True)]], line_spacing=1.3)

# =====================================================================
# P7 技术底座：三台计算机
# =====================================================================
s = new_slide()
header(s, "TECH FOUNDATION", "技术底座：支撑物理 AI 的“三台计算机”", 7)
deco_circles(s)
comps = [
    ("训练计算机", "超级算力“教练”", "DGX / 云 GPU 集群", "用海量数据与世界模型训练出机器人的“大脑”", TEAL, MINT_BG),
    ("仿真计算机", "虚拟练兵场", "Omniverse + Cosmos", "数字孪生 + 世界模型，生成合成数据并做闭环验证", SKY, SKY_BG),
    ("本体计算机", "装在机器里的小脑", "Jetson Thor / DRIVE AGX", "部署在机器人与车辆上，实时感知、决策、控制", AMBER, AMBER_BG),
]
for i, (t, role, prod, d, c, bgc) in enumerate(comps):
    x = MARGIN + i * Inches(4.13)
    card(s, x, Inches(1.75), Inches(3.85), Inches(3.1), fill=bgc, line=None)
    add_shape(s, MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.05), Inches(0.7), Inches(0.7),
              fill=WHITE, line=c, line_w=1.5)
    add_text(s, x + Inches(0.3), Inches(2.22), Inches(0.7), Inches(0.4),
             [[(str(i + 1), 18, c, True)]], align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.15), Inches(2.1), Inches(2.6), Inches(0.4), [[(t, 17, INK, True)]])
    add_text(s, x + Inches(1.15), Inches(2.5), Inches(2.6), Inches(0.3), [[(role, 11, MUTED, False)]])
    hline(s, x + Inches(0.3), Inches(2.95), Inches(3.25), color=WHITE, weight=2.0)
    add_text(s, x + Inches(0.3), Inches(3.1), Inches(3.3), Inches(0.35),
             [[("代表平台：" + prod, 11.5, c, True)]])
    add_text(s, x + Inches(0.3), Inches(3.5), Inches(3.3), Inches(1.2),
             [[(d, 12, BODY, False)]], line_spacing=1.3)
    if i < 2:
        add_shape(s, MSO_SHAPE.CHEVRON, x + Inches(3.87), Inches(3.1), Inches(0.3), Inches(0.4), fill=LINE)
card(s, MARGIN, Inches(5.2), Inches(12.1), Inches(1.55), fill=WHITE, line=LINE)
add_text(s, MARGIN + Inches(0.4), Inches(5.42), Inches(11.4), Inches(0.4),
         [[("关键模型层", 14, INK, True)]])
models = [("Cosmos 世界模型", "物理世界的“教科书”，基于 2000 万+ 小时真实数据训练"),
          ("GR00T 机器人模型", "通用人形机器人的“大脑”，开源可用"),
          ("Alpamayo 自驾模型", "面向 L4 自动驾驶的开源推理模型")]
for i, (t, d) in enumerate(models):
    x = MARGIN + Inches(0.4) + i * Inches(3.85)
    add_text(s, x, Inches(5.85), Inches(3.7), Inches(0.35), [[(t, 12.5, TEAL_DK, True)]])
    add_text(s, x, Inches(6.18), Inches(3.7), Inches(0.55), [[(d, 11, BODY, False)]], line_spacing=1.2)

# =====================================================================
# P8 章节页 02
# =====================================================================
section_slide("02", "与 Agent、机器人的对比", "VS. AGENTS & ROBOTS",
              ["物理 AI vs AI Agent：数字世界 vs 物理世界",
               "物理 AI vs 传统机器人：专机专用 vs 通用智能",
               "物理 AI 的四大核心优势"], 8)

# =====================================================================
# P9 物理 AI vs AI Agent
# =====================================================================
s = new_slide()
header(s, "COMPARISON 1", "物理 AI vs AI Agent：一字之差，两个世界", 9)
deco_circles(s)
rows = [
    ("运行空间", "数字世界：浏览器、软件、API", "物理世界：工厂、道路、家庭"),
    ("输出结果", "文字、代码、点击等数字动作", "电机指令：抓取、行走、驾驶"),
    ("核心难点", "推理规划、工具调用、长任务", "物理规律、实时控制、安全冗余"),
    ("试错成本", "错了重来，几乎零成本", "试错昂贵危险，依赖仿真训练"),
    ("典型代表", "Claude / Manus / 各类办公 Agent", "人形机器人、自动驾驶汽车"),
]
tx = [MARGIN + Inches(2.5), MARGIN + Inches(2.5) + Inches(4.85)]
colw = Inches(4.65)
# 表头
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, tx[0], Inches(1.6), colw, Inches(0.55), fill=VIOLET_BG, adj=0.25)
add_text(s, tx[0], Inches(1.72), colw, Inches(0.35), [[("AI Agent（代理 AI）", 15, VIOLET, True)]], align=PP_ALIGN.CENTER)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, tx[1], Inches(1.6), colw, Inches(0.55), fill=MINT, adj=0.25)
add_text(s, tx[1], Inches(1.72), colw, Inches(0.35), [[("物理 AI", 15, TEAL_DK, True)]], align=PP_ALIGN.CENTER)
add_shape(s, MSO_SHAPE.OVAL, tx[0] + colw + Inches(0.03), Inches(1.68), Inches(0.38), Inches(0.38),
          fill=WHITE, line=MUTED, line_w=1.0)
add_text(s, tx[0] + colw + Inches(0.03), Inches(1.77), Inches(0.38), Inches(0.24),
         [[("VS", 10, BODY, True)]], align=PP_ALIGN.CENTER)
yy = 2.32
for i, (dim, a, b) in enumerate(rows):
    rh = 0.78
    add_text(s, MARGIN + Inches(0.1), Inches(yy + 0.16), Inches(2.2), Inches(0.4),
             [[(dim, 13.5, INK, True)]])
    card(s, tx[0], Inches(yy), colw, Inches(rh - 0.12), fill=SOFT, line=None, radius=0.18)
    add_text(s, tx[0] + Inches(0.2), Inches(yy + 0.12), colw - Inches(0.4), Inches(0.5),
             [[(a, 12, BODY, False)]], line_spacing=1.1)
    card(s, tx[1], Inches(yy), colw, Inches(rh - 0.12), fill=MINT_BG, line=None, radius=0.18)
    add_text(s, tx[1] + Inches(0.2), Inches(yy + 0.12), colw - Inches(0.4), Inches(0.5),
             [[(b, 12, TEAL_DK, True)]], line_spacing=1.1)
    yy += rh
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.35), Inches(12.1), Inches(0.6),
          fill=WHITE, line=LINE, adj=0.5)
add_text(s, MARGIN, Inches(6.49), Inches(12.1), Inches(0.4),
         [[("关系：Agent 是物理 AI 的“决策层”——物理 AI = Agent 的大脑 + 机器人的身体 + 物理常识。", 13.5, INK, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P10 物理 AI vs 传统机器人
# =====================================================================
s = new_slide()
header(s, "COMPARISON 2", "物理 AI vs 传统机器人：从“专机”到“通才”", 10)
deco_circles(s)
rows = [
    ("工作方式", "工程师逐行写死程序", "大模型理解指令，自主规划动作"),
    ("任务范围", "单一任务：焊接、拧螺丝", "一机多能：搬箱、分拣、装配都会"),
    ("环境适应", "环境必须精确布置，怕变化", "能适应杂乱、动态的真实场景"),
    ("学习进化", "换任务 = 重新编程调试数周", "换任务 = 给数据/演示，快速泛化"),
    ("典型形态", "工业机械臂、AGV 小车", "人形机器人、四足、智能车辆"),
]
tx = [MARGIN + Inches(2.5), MARGIN + Inches(2.5) + Inches(4.85)]
colw = Inches(4.65)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, tx[0], Inches(1.6), colw, Inches(0.55), fill=SOFT, adj=0.25)
add_text(s, tx[0], Inches(1.72), colw, Inches(0.35), [[("传统机器人", 15, BODY, True)]], align=PP_ALIGN.CENTER)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, tx[1], Inches(1.6), colw, Inches(0.55), fill=MINT, adj=0.25)
add_text(s, tx[1], Inches(1.72), colw, Inches(0.35), [[("物理 AI 驱动的机器人", 15, TEAL_DK, True)]], align=PP_ALIGN.CENTER)
add_shape(s, MSO_SHAPE.OVAL, tx[0] + colw + Inches(0.03), Inches(1.68), Inches(0.38), Inches(0.38),
          fill=WHITE, line=MUTED, line_w=1.0)
add_text(s, tx[0] + colw + Inches(0.03), Inches(1.77), Inches(0.38), Inches(0.24),
         [[("VS", 10, BODY, True)]], align=PP_ALIGN.CENTER)
yy = 2.32
for dim, a, b in rows:
    rh = 0.78
    add_text(s, MARGIN + Inches(0.1), Inches(yy + 0.16), Inches(2.2), Inches(0.4),
             [[(dim, 13.5, INK, True)]])
    card(s, tx[0], Inches(yy), colw, Inches(rh - 0.12), fill=SOFT, line=None, radius=0.18)
    add_text(s, tx[0] + Inches(0.2), Inches(yy + 0.12), colw - Inches(0.4), Inches(0.5),
             [[(a, 12, BODY, False)]], line_spacing=1.1)
    card(s, tx[1], Inches(yy), colw, Inches(rh - 0.12), fill=MINT_BG, line=None, radius=0.18)
    add_text(s, tx[1] + Inches(0.2), Inches(yy + 0.12), colw - Inches(0.4), Inches(0.5),
             [[(b, 12, TEAL_DK, True)]], line_spacing=1.1)
    yy += rh
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.35), Inches(12.1), Inches(0.6),
          fill=WHITE, line=LINE, adj=0.5)
add_text(s, MARGIN, Inches(6.49), Inches(12.1), Inches(0.4),
         [[("本质变化：机器人从“自动化设备”升级为“可自主学习的智能体”——这叫“专业通才型”机器人。", 13.5, INK, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P11 物理 AI 的四大优势
# =====================================================================
s = new_slide()
header(s, "ADVANTAGES", "物理 AI 的核心优势", 11)
deco_circles(s)
advs = [
    ("通用性", "一套模型适配多种本体与任务，能力可迁移、可复制，不再“一个场景一套系统”", "A", TEAL, MINT_BG),
    ("训练快、成本低", "仿真中并行训练，一夜完成百万次试错；合成数据把采集成本降低几个数量级", "B", SKY, SKY_BG),
    ("安全可控", "危险场景先在虚拟世界验证，闭环测试充分后才部署到真机，大幅降低事故风险", "C", AMBER, AMBER_BG),
    ("持续进化", "真实运行数据回流再训练，形成“越用越聪明”的数据飞轮", "D", VIOLET, VIOLET_BG),
]
for i, (t, d, tag, c, bgc) in enumerate(advs):
    x = MARGIN + (i % 2) * Inches(6.15)
    y = Inches(1.75) + (i // 2) * Inches(2.4)
    card(s, x, y, Inches(5.85), Inches(2.15), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.OVAL, x + Inches(0.35), y + Inches(0.35), Inches(0.85), Inches(0.85), fill=bgc)
    add_text(s, x + Inches(0.35), y + Inches(0.55), Inches(0.85), Inches(0.45),
             [[(tag, 20, c, True)]], align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.45), y + Inches(0.32), Inches(4.2), Inches(0.45),
             [[(t, 18, INK, True)]])
    add_text(s, x + Inches(1.45), y + Inches(0.85), Inches(4.15), Inches(1.1),
             [[(d, 12.5, BODY, False)]], line_spacing=1.35)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.35), Inches(12.1), Inches(0.6),
          fill=MINT, line=None, adj=0.5)
add_text(s, MARGIN, Inches(6.49), Inches(12.1), Inches(0.4),
         [[("一句话：物理 AI 让“造一个什么都能干的机器人”从科幻变成了工程问题。", 14, TEAL_DK, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P12 章节页 03
# =====================================================================
section_slide("03", "产业格局与优秀案例", "LANDSCAPE & CASES",
              ["国际领先公司与平台型玩家",
               "中国力量：量产与场景落地",
               "六个值得关注的标杆案例"], 12)

# =====================================================================
# P13 国际公司
# =====================================================================
s = new_slide()
header(s, "GLOBAL PLAYERS", "哪些公司做得好 · 国际篇", 13)
deco_circles(s)
intl = [
    ("NVIDIA 英伟达", "平台底座", "全栈物理 AI 平台：Cosmos 世界模型、GR00T、Omniverse 仿真、Thor 芯片，生态最完整", TEAL),
    ("Tesla 特斯拉", "整机+自驾", "Optimus V3 产线 2026 年中启动，远期规划百万台/年；FSD 自动驾驶数据闭环领先", SKY),
    ("Figure AI", "人形机器人", "BotQ 工厂每小时下线一台 Figure 03，已进入宝马工厂承担真实生产物流", AMBER),
    ("Physical Intelligence", "具身大模型", "π 系列通用机器人基础模型，让不同机器人共享一个“大脑”", VIOLET),
    ("Google DeepMind", "模型研发", "Gemini Robotics 系列：把多模态大模型能力延伸到机器人操作", ROSE),
    ("Boston Dynamics", "运动控制", "Atlas 电动版与现代汽车工厂试点结合，运动控制能力业界标杆", TEAL_DK),
]
for i, (t, tag, d, c) in enumerate(intl):
    x = MARGIN + (i % 2) * Inches(6.15)
    y = Inches(1.7) + (i // 2) * Inches(1.72)
    card(s, x, y, Inches(5.85), Inches(1.52), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.28), y + Inches(0.24), Inches(1.5), Inches(0.34),
              fill=SOFT, adj=0.5)
    add_text(s, x + Inches(0.28), y + Inches(0.3), Inches(1.5), Inches(0.26),
             [[(tag, 10.5, c, True)]], align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.95), y + Inches(0.24), Inches(3.8), Inches(0.4),
             [[(t, 15, INK, True)]])
    add_text(s, x + Inches(0.28), y + Inches(0.72), Inches(5.35), Inches(0.7),
             [[(d, 11.5, BODY, False)]], line_spacing=1.25)

# =====================================================================
# P14 中国公司
# =====================================================================
s = new_slide()
header(s, "CHINA PLAYERS", "哪些公司做得好 · 中国篇", 14)
deco_circles(s)
add_text(s, MARGIN, Inches(1.4), Inches(12.1), Inches(0.4),
         [[("2026 年被视为人形机器人“量产元年”：中国整机产量有望突破 10 万台，出货量占全球近九成", 13, AMBER, True)]])
chn = [
    ("宇树科技 Unitree", "科创板 IPO 已注册生效，双足人形出货量全球第一；2026 年出货目标 1-2 万台，完成全球首例人形机器人手术演示", TEAL),
    ("智元机器人 AgiBot", "2025 年出货 5168 台全球第一（Omdia），第 1.5 万台已下线；开源 AgiBot World 百万条轨迹数据集与 GO-1 模型", SKY),
    ("银河通用 Galbot", "融资 25 亿元、估值超 200 亿；主打仿真合成数据训练，机器人落地零售、工厂、药房场景", AMBER),
    ("优必选 UBTech", "Walker S2 工业人形机器人获超 13 亿元订单，批量进入汽车工厂实训", VIOLET),
    ("星动纪元 / 逐际动力", "清华系创业代表：端到端具身大模型 + 高动态运动控制，科研与工业场景双线推进", ROSE),
    ("自动驾驶阵营", "华为、 Momenta、小鹏、文远知行等推动 L2++ 普及与 L4 试运营，是物理 AI 最成熟的落地场景", TEAL_DK),
]
for i, (t, d, c) in enumerate(chn):
    x = MARGIN + (i % 2) * Inches(6.15)
    y = Inches(1.95) + (i // 2) * Inches(1.68)
    card(s, x, y, Inches(5.85), Inches(1.48), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.12), Inches(1.48), fill=c, adj=0.5)
    add_text(s, x + Inches(0.32), y + Inches(0.16), Inches(5.3), Inches(0.4),
             [[(t, 14.5, INK, True)]])
    add_text(s, x + Inches(0.32), y + Inches(0.58), Inches(5.35), Inches(0.85),
             [[(d, 11, BODY, False)]], line_spacing=1.22)

# =====================================================================
# P15 优秀案例
# =====================================================================
s = new_slide()
header(s, "SHOWCASES", "值得关注的优秀案例", 15)
deco_circles(s)
cases = [
    ("Figure 03 进宝马工厂", "2026 年 6 月起在宝马斯帕坦堡工厂承担真实生产物流，人形机器人首次大规模“进厂打工”", "制造", TEAL, MINT_BG),
    ("宇树机器人手术演示", "完成全球首例人形机器人辅助活体手术演示，展示亚毫米级操作精度潜力", "医疗", ROSE, ROSE_BG),
    ("奔驰 CLA 搭载 Alpamayo", "首款基于 NVIDIA 开源推理模型的量产乘用车，向 L4 级自动驾驶迈进", "出行", SKY, SKY_BG),
    ("普智 G2 零失误实战", "在国内头部 3C 工厂完成 8 小时、2283 次任务零失误测试，验证工业可靠性", "制造", AMBER, AMBER_BG),
    ("智元 AgiBot World", "开源 100 万+ 条真机操作轨迹，覆盖 217 项任务，成为行业公共数据底座", "开源", VIOLET, VIOLET_BG),
    ("银河通用智慧零售", "机器人在无人药店、便利店完成取货补货全流程，7×24 小时商业化运营", "零售", TEAL_DK, MINT_BG),
]
for i, (t, d, tag, c, bgc) in enumerate(cases):
    x = MARGIN + (i % 3) * Inches(4.13)
    y = Inches(1.7) + (i // 3) * Inches(2.6)
    card(s, x, y, Inches(3.85), Inches(2.35), fill=bgc, line=None)
    chip(s, x + Inches(0.28), y + Inches(0.25), Inches(0.85), tag, fg=WHITE, bg=c)
    add_text(s, x + Inches(0.28), y + Inches(0.72), Inches(3.3), Inches(0.65),
             [[(t, 14.5, INK, True)]], line_spacing=1.15)
    add_text(s, x + Inches(0.28), y + Inches(1.32), Inches(3.3), Inches(0.95),
             [[(d, 11, BODY, False)]], line_spacing=1.25)

# =====================================================================
# P16 章节页 04
# =====================================================================
section_slide("04", "未来趋势与挑战", "TRENDS & CHALLENGES",
              ["五大发展趋势：量产、世界模型、VLA、数据飞轮、成本",
               "四大现实挑战：安全、数据、灵巧操作、法规伦理"], 16)

# =====================================================================
# P17 五大趋势
# =====================================================================
s = new_slide()
header(s, "TRENDS", "物理 AI 的五大发展趋势", 17)
deco_circles(s)
trends = [
    ("量产元年开启", "2026 年中国人形机器人产量有望破 10 万台；Figure、特斯拉、智元、宇树纷纷建产线，规模效应显现", TEAL),
    ("世界模型成为标配", "Cosmos 类模型让 AI 在“脑海”中预演物理世界，成为训练与决策的基础设施", SKY),
    ("VLA 大模型驱动", "视觉-语言-动作模型让机器人“听懂人话、看懂场景、直接动手”，走向通用", VIOLET),
    ("数据飞轮加速", "量产越多→真实数据越多→模型越强→越好卖，头部效应将愈发明显", AMBER),
    ("成本快速下降", "整机价格进入 10 万元级甚至更低，租赁模式兴起，中小企业与个人用得起", ROSE),
]
yy = 1.66
for i, (t, d, c) in enumerate(trends):
    card(s, MARGIN, Inches(yy), Inches(12.1), Inches(0.92), fill=WHITE, line=LINE, radius=0.14)
    add_shape(s, MSO_SHAPE.OVAL, MARGIN + Inches(0.25), Inches(yy + 0.19), Inches(0.54), Inches(0.54), fill=MINT)
    add_text(s, MARGIN + Inches(0.25), Inches(yy + 0.31), Inches(0.54), Inches(0.3),
             [[(f"{i+1}", 15, TEAL_DK, True)]], align=PP_ALIGN.CENTER)
    add_text(s, MARGIN + Inches(1.0), Inches(yy + 0.13), Inches(2.85), Inches(0.7),
             [[(t, 15.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN + Inches(4.05), Inches(yy + 0.13), Inches(7.8), Inches(0.7),
             [[(d, 12, BODY, False)]], line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)
    yy += 1.04
add_text(s, MARGIN, Inches(6.95 - 0.35), Inches(12.1), Inches(0.35),
         [[("业界共识：未来 3-5 年单一场景任务基本可行，2030 年前后进入多场景规模化期。", 12.5, TEAL_DK, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P18 挑战与风险
# =====================================================================
s = new_slide()
header(s, "CHALLENGES", "也要看到：四大现实挑战", 18)
deco_circles(s)
chls = [
    ("安全与可靠", "几十公斤的机器人与人共处，一次失控就是事故；需要功能安全、冗余设计与行业标准兜底", "01"),
    ("数据仍然稀缺", "真实操作数据采集慢、贵；仿真与现实的差距（Sim-to-Real Gap）仍需持续弥合", "02"),
    ("灵巧操作瓶颈", "行走已较成熟，但“手”的精细操作（柔性物体、小零件）仍是世界级难题", "03"),
    ("成本与法规", "整机、维护成本仍高；责任认定、隐私、就业影响等伦理法规框架尚在早期", "04"),
]
for i, (t, d, no) in enumerate(chls):
    x = MARGIN + (i % 2) * Inches(6.15)
    y = Inches(1.75) + (i // 2) * Inches(2.1)
    card(s, x, y, Inches(5.85), Inches(1.85), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.12), Inches(1.85), fill=ROSE, adj=0.5)
    add_text(s, x + Inches(0.35), y + Inches(0.2), Inches(4.6), Inches(0.4),
             [[(t, 16, INK, True)]])
    add_text(s, x + Inches(4.7), y + Inches(0.14), Inches(1.0), Inches(0.6),
             [[(no, 26, LINE, True)]], align=PP_ALIGN.RIGHT)
    add_text(s, x + Inches(0.35), y + Inches(0.7), Inches(5.25), Inches(1.0),
             [[(d, 12, BODY, False)]], line_spacing=1.3)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.15), Inches(12.1), Inches(0.75),
          fill=ROSE_BG, line=None, adj=0.2)
add_text(s, MARGIN + Inches(0.4), Inches(6.32), Inches(11.5), Inches(0.45),
         [[("理性看待：物理 AI 是一条前景广阔但更漫长的路，需要定力与耐心，避免短期过度乐观。", 13.5, ROSE, True)]],
         align=PP_ALIGN.CENTER)

# =====================================================================
# P19 普通人如何学习
# =====================================================================
s = new_slide()
header(s, "HOW TO LEARN", "普通人如何入门物理 AI？", 19)
deco_circles(s)
steps = [
    ("第 1 步 · 建立认知", "看懂概念框架", ["关注 NVIDIA GTC / CES 演讲、行业报告", "理解“三台计算机”与仿真优先思路"]),
    ("第 2 步 · 动手仿真", "零硬件成本起步", ["NVIDIA Isaac Sim / Isaac Lab 免费教程", "MuJoCo、Gazebo 仿真器 + ROS 2 基础"]),
    ("第 3 步 · 玩转模型", "用开源资源实操", ["Hugging Face LeRobot 低门槛入门", "智元 AgiBot World 等开源数据集"]),
    ("第 4 步 · 真机实践", "加入社区共创", ["万元级开源机械臂/小车 DIY", "参加机器人竞赛、黑客松与开发者社区"]),
]
for i, (t, st, pts) in enumerate(steps):
    x = MARGIN + i * Inches(3.08)
    card(s, x, Inches(1.7), Inches(2.85), Inches(3.4), fill=WHITE, line=LINE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.85), Inches(0.1), fill=TEAL, adj=0.5)
    add_text(s, x + Inches(0.22), Inches(1.95), Inches(2.4), Inches(0.4), [[(t, 13.5, INK, True)]])
    add_text(s, x + Inches(0.22), Inches(2.36), Inches(2.4), Inches(0.3), [[(st, 11, TEAL_DK, True)]])
    hline(s, x + Inches(0.22), Inches(2.75), Inches(2.4), color=LINE, weight=1.0)
    yy = 2.9
    for p in pts:
        add_shape(s, MSO_SHAPE.OVAL, x + Inches(0.24), Inches(yy + 0.07), Inches(0.08), Inches(0.08), fill=TEAL)
        add_text(s, x + Inches(0.42), Inches(yy), Inches(2.3), Inches(0.95),
                 [[(p, 11, BODY, False)]], line_spacing=1.18)
        yy += 1.0
    if i < 3:
        add_shape(s, MSO_SHAPE.CHEVRON, x + Inches(2.87), Inches(3.1), Inches(0.24), Inches(0.36), fill=MINT)
card(s, MARGIN, Inches(5.4), Inches(12.1), Inches(1.4), fill=MINT_BG, line=None)
add_text(s, MARGIN + Inches(0.4), Inches(5.58), Inches(11.4), Inches(0.4),
         [[("给非技术同学的建议", 14, TEAL_DK, True)]])
add_text(s, MARGIN + Inches(0.4), Inches(5.95), Inches(11.4), Inches(0.75),
         [[("不必先学编程：从理解产业逻辑、体验机器人产品、跟踪头部公司动态开始；", 12.5, BODY, False)],
          [("把物理 AI 当作“新基建”来理解 —— 它带来的岗位与机会远不止算法工程师。", 12.5, BODY, False)]],
         line_spacing=1.3)

# =====================================================================
# P20 总结
# =====================================================================
s = new_slide()
header(s, "SUMMARY", "总结：物理 AI 离我们有多近？", 20)
deco_circles(s)
sums = [
    ("它是什么", "能理解物理世界并与之交互的 AI，是感知/生成/代理 AI 之后的第四波浪潮", TEAL, MINT_BG),
    ("它在做什么", "用仿真与世界模型低成本训练，再部署到机器人与自动驾驶等自主机器", SKY, SKY_BG),
    ("为什么重要", "让机器人从“专机”变“通才”，开启新一轮工业革命的入口", AMBER, AMBER_BG),
]
for i, (t, d, c, bgc) in enumerate(sums):
    x = MARGIN + i * Inches(4.13)
    card(s, x, Inches(1.7), Inches(3.85), Inches(1.95), fill=bgc, line=None)
    add_text(s, x + Inches(0.3), Inches(1.95), Inches(3.2), Inches(0.4), [[(t, 16, INK, True)]])
    add_text(s, x + Inches(0.3), Inches(2.42), Inches(3.3), Inches(1.1),
             [[(d, 12, BODY, False)]], line_spacing=1.3)
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(4.0), Inches(12.1), Inches(1.5),
          fill=MINT, line=None, adj=0.14)
add_text(s, MARGIN + Inches(0.5), Inches(4.28), Inches(11.1), Inches(1.0),
         [[("“物理 AI 与机器人技术，将开启新一轮工业革命。” —— 黄仁勋", 17, TEAL_DK, True)],
          [("对我们而言：现在正是理解它、用好它的最佳时点。", 14, INK, False)]],
         line_spacing=1.5, align=PP_ALIGN.CENTER)
add_text(s, MARGIN, Inches(5.95), Inches(12.1), Inches(0.6),
         [[("谢谢聆听 · Q&A", 26, INK, True)]], align=PP_ALIGN.CENTER)
add_text(s, MARGIN, Inches(6.6), Inches(12.1), Inches(0.35),
         [[("资料来源：NVIDIA 官方资料、CES 2026、新华网、Omdia、36 氪等公开报道（2026 年 7 月）", 10, MUTED, False)]],
         align=PP_ALIGN.CENTER)

prs.save("/workspace/物理AI-科普分享.pptx")
print("OK, slides:", len(prs.slides._sldIdLst))
